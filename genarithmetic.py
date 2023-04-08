#!/usr/bin/env python3

import getopt
import pptx
import random
import sys

def setTextBox (txBox, text):
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Calibri'
    font.size = pptx.util.Pt(17)
    font.color.rgb = pptx.dml.color.RGBColor (0x00, 0x00, 0x00)

def genSlide (presentation, texts):
    # 6 is blank slide layout
    blankSlideLayout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide (blankSlideLayout)

    width = pptx.util.Inches (2)
    height = pptx.util.Inches (10)
    top = pptx.util.Inches (0.5)

    # left column
    left = pptx.util.Inches (0.5)

    txBox = slide.shapes.add_textbox (left, top, width, height)
    setTextBox (txBox, texts[0])

    # middle column
    left = pptx.util.Inches (3)

    txBox = slide.shapes.add_textbox (left, top, width, height)
    setTextBox (txBox, texts[1])

    # right column
    left = pptx.util.Inches (5.5)
    txBox = slide.shapes.add_textbox (left, top, width, height)
    setTextBox (txBox, texts[2])

def createPresentation ():
    presentation = pptx.Presentation ()

    # Change the slide to 8.5 x 11 inches portrait
    presentation.slide_width = pptx.util.Inches (8.5)
    presentation.slide_height = pptx.util.Inches (11)
    return presentation

def getProblemString (a, b, c, d, answer):
    if answer:
        solution = a * b + c * d
        if (c < 0):
            return ("%d × %d - %d = %d\n" % (a, b, d, solution))
        else:
            return ("%d × %d + %d = %d\n" % (a, b, d, solution))
    else:
        if (c < 0):
            return ("%d × %d - %d =\n" % (a, b, d))
        else:
            return ("%d × %d + %d =\n" % (a, b, d))

def getMaths (num):
    problem = ''
    solution = ''

    n = 0
    while n < num:
        a = random.randrange (2, 11)
        b = random.randrange (2, 11)
        c = random.randrange (0, 2) * 2 - 1;
        d = random.randrange (0, 100)

        if (a * b + c * d) < 0:
            continue

        problem += getProblemString (a, b, c, d, False)
        solution += getProblemString (a, b, c, d, True)

        n += 1

    return (problem, solution)

def usage (prog):
    text = """usage: %s [options]
options:
    -h,--help           displays this help information, then exit
    -p,--page number    number of pages to be generated.  default is 1
"""
    print (text % prog)

def main (argv):
    prog = argv[0]
    numPages = 1

    try:
        opts, args = getopt.getopt(sys.argv[1:], "hp:", ["help", "pages"])
    except getopt.GetoptError as err:
        print (err)
        usage(prog)
        exit (2)

    for o, a in opts:
        if o in ('-h', '--help'):
            usage (prog)
            exit (0)
        elif o in ('-p', '--pages'):
            try:
                numPages = int (a)
            except:
                print ('invalid page value')
                exit (1)
            if numPages < 1:
                print ('invalid page value')
                exit (1)

    random.seed ()

    problemPresentation = createPresentation ()
    solutionPresentation = createPresentation ()

    for page in range (numPages):
        problemTexts = []
        solutionTexts = []

        p, s = getMaths (33)
        problemTexts.append (p)
        solutionTexts.append (s)

        p, s = getMaths (33)
        problemTexts.append (p)
        solutionTexts.append (s)

        p, s = getMaths (34)
        problemTexts.append (p)
        solutionTexts.append (s)

        genSlide (problemPresentation, problemTexts)
        genSlide (solutionPresentation, solutionTexts)

    problemPresentation.save ("problem.pptx")
    solutionPresentation.save ("solution.pptx")

if __name__ == "__main__":
    main (sys.argv)
